import os
import json
import fitz  # PyMuPDF
import sys
from pathlib import Path
from tkinter import filedialog, Tk, simpledialog, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_openai import OpenAI as LangchainOpenAI

# Langchain LLM 객체 생성
llm = LangchainOpenAI(
    openai_api_base="http://localhost:1234/v1",
    openai_api_key="lm-studio",
    model_name="local-model",
    temperature=0.7,
)

MAX_EXAMPLES = 5  # 프롬프트 길이 방지를 위한 최대 예시 수

def extract_text_and_images(pdf_path, img_output_dir="extracted_images"):
    doc = fitz.open(pdf_path)
    os.makedirs(img_output_dir, exist_ok=True)
    full_text = ""
    page_images = {}

    for i, page in enumerate(doc):
        full_text += page.get_text()
        pix = page.get_pixmap(dpi=300)
        img_path = os.path.join(img_output_dir, f"page_{i+1}.png")
        pix.save(img_path)
        page_images[i+1] = [img_path]

    return full_text, img_output_dir, page_images

def split_text_to_slides(text, max_chars=800):
    paragraphs = text.split('\n')
    slides = []
    current_slide = ""

    for para in paragraphs:
        if len(current_slide) + len(para) + 1 < max_chars:
            current_slide += para + '\n'
        else:
            slides.append(current_slide.strip())
            current_slide = para + '\n'
    if current_slide:
        slides.append(current_slide.strip())

    return slides

def match_image_by_index(index, page_images):
    page_num = index + 1
    return page_images.get(page_num, [None])[0]

def ask_llm_with_examples(examples, new_question):
    example_prompt = ""
    for ex in examples:
        example_prompt += f"Q: {ex['question']}\nA: {ex['answer']}\n\n"

    prompt = f"""다음은 브로셔 내용을 바탕으로 한 예시 Q&A입니다:

    {example_prompt}
    Q: {new_question}
    A:"""

    estimated_tokens = len(prompt.split())
    print("=" * 50)
    print(f"📤 LLM 호출 프롬프트 (예상 토큰 수: {estimated_tokens})")
    print(prompt)
    print("=" * 50)

    try:
        response = llm.invoke(prompt)
        if not response.strip():
            print("⚠️ LLM이 빈 응답을 반환했습니다.")
            return "(LLM 응답 없음)"
        print("📥 LLM 응답:\n", response)
        return response.strip()
    except Exception as e:
        print(f"❌ LLM 호출 실패: {e}")
        return f"(LLM 호출 오류: {e})"

def create_ppt(slide_contents, page_images, output_filename):
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "목차"
    content = title_slide.placeholders[1].text_frame

    for idx, (title, _) in enumerate(slide_contents, 1):
        content.add_paragraph().text = f"{idx}. {title}"

    for index, (title, content_text) in enumerate(slide_contents):
        chunks = split_text_to_slides(content_text)
        for i, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide_title = title if i == 0 else f"{title} (계속)"
            slide.shapes.title.text = slide_title

            para = slide.shapes.title.text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
                run.font.size = Pt(18)
                run.font.name = "맑은 고딕"

            content_box = slide.placeholders[1]
            content_box.text = chunk
            for p in content_box.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.name = "맑은 고딕"

            matched_image = match_image_by_index(index, page_images)
            if matched_image:
                slide.shapes.add_picture(matched_image, Inches(5.5), Inches(1.5), width=Inches(3))

    prs.save(output_filename)

def parse_prompt_template(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()
        start = content.find("[")
        end = content.rfind("]") + 1
        try:
            return json.loads(content[start:end])
        except json.JSONDecodeError:
            messagebox.showerror("오류", "템플릿 파일의 JSON 형식이 잘못되었습니다.")
            raise

def main():
    root = Tk()
    root.withdraw()

    pdf_path = filedialog.askopenfilename(title="PDF 파일 선택", filetypes=[("PDF files", "*.pdf")])
    if not pdf_path:
        messagebox.showwarning("경고", "PDF 파일을 선택하지 않았습니다.")
        return

    template_path = filedialog.askopenfilename(title="프롬프트 템플릿 파일 선택", filetypes=[("Python files", "*.py")])
    if not template_path:
        messagebox.showwarning("경고", "템플릿 파일을 선택하지 않았습니다.")
        return

    prompt_templates = parse_prompt_template(template_path)

    num_examples = simpledialog.askinteger("예시 개수", f"Few-shot 예시 개수를 입력하세요 (최대 {MAX_EXAMPLES}개):", minvalue=1, maxvalue=len(prompt_templates))
    if num_examples is None:
        messagebox.showwarning("경고", "예시 개수를 입력하지 않았습니다.")
        return
    if num_examples > MAX_EXAMPLES:
        messagebox.showinfo("알림", f"예시 개수가 많아 {MAX_EXAMPLES}개로 제한합니다.")
        num_examples = MAX_EXAMPLES

    output_path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                               filetypes=[("PowerPoint files", "*.pptx")],
                                               title="저장할 PPT 파일 이름")
    if not output_path:
        messagebox.showwarning("경고", "저장 파일명을 입력하지 않았습니다.")
        return

    print("📄 문서 분석 중...")
    text, img_dir, page_images = extract_text_and_images(pdf_path)

    print("🧠 LLM을 통한 답변 생성 중...")
    slides = []
    few_shot_examples = prompt_templates[:num_examples]
    generate_targets = prompt_templates

    for template in generate_targets:
        title = template.get("title", "제목 없음")
        question = template.get("question", "")
        print(f"📌 '{title}' 새 답변 생성 중...")
        llm_answer = ask_llm_with_examples(few_shot_examples, question)
        original_answer = template.get("answer", "(없음)")
        #comparison = f"[LLM 응답]\n{llm_answer}\n\n[기존 답변]\n{original_answer}"
        comparison = f"{llm_answer}"
        slides.append((title, comparison))

    print("🎞 PPT 생성 중...")
    create_ppt(slides, page_images, output_path)
    print("✅ 완료! 저장 위치:", output_path)

    # 루트 윈도우를 앞으로 가져오고 메시지박스를 맨 위에 표시
    root.deiconify()
    root.lift()
    root.attributes("-topmost", True)
    root.after(100, lambda: root.attributes("-topmost", False))
    messagebox.showinfo("완료", f"PPT 생성이 완료되었습니다:\n{output_path}")

    root.destroy()


if __name__ == "__main__":
    main()
